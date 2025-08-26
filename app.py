# app.py
# Text → PowerPoint (Template-Aware) — Streamlit app
# Public MVP supporting OpenAI / Anthropic / Gemini / OpenRouter / Custom OpenAI-compatible via user-supplied API key
# No AI image generation; reuses images embedded in the uploaded template (.pptx/.potx)
# MIT License

import io
import json
import re
import zipfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

APP_TITLE = "Text → PowerPoint (Template-Aware)"
MAX_UPLOAD_MB = 25

# --------------------- Utility: Requests session with retries ---------------------

def make_session():
    s = requests.Session()
    retries = Retry(
        total=3,
        backoff_factor=0.5,
        status_forcelist=[429, 500, 502, 503, 504],
        allowed_methods=["POST", "GET"],
        raise_on_status=False,
    )
    adapter = HTTPAdapter(max_retries=retries)
    s.mount("https://", adapter)
    s.mount("http://", adapter)
    return s

session = make_session()

# --------------------- Utility: JSON extraction ---------------------

def extract_json(text: str) -> Dict[str, Any]:
    """Try to parse JSON; if that fails, greedily extract the first JSON object or array."""
    if not text:
        return {"slides": []}
    try:
        return json.loads(text)
    except Exception:
        pass
    m = re.search(r"\{[\s\S]*\}", text)
    if m:
        try:
            return json.loads(m.group(0))
        except Exception:
            pass
    m = re.search(r"\[[\s\S]*\]", text)
    if m:
        try:
            return {"slides": json.loads(m.group(0))}
        except Exception:
            pass
    return {"slides": [{"title": "Overview", "bullets": [text[:2000]]}]}

# --------------------- LLM Clients ---------------------

LLM_MODELS = {
    "OpenAI": "gpt-4o-mini",
    "Anthropic": "claude-3-5-sonnet-20240620",
    "Gemini": "gemini-1.5-pro",
    "OpenRouter": "openrouter/auto",
    "Custom (OpenAI-compatible)": "gpt-4o-mini",  # user can override
}

SYSTEM_PROMPT = (
    "You are a senior presentation designer. Turn input text into a clean slide deck plan. "
    "Return only JSON matching this schema: {\n"
    '  "title": string (optional),\n'
    '  "slides": [\n'
    "    {\n"
    "      \"layout_hint\": one of ['title','section','title+content','two-content','quote','bullets-only','image-left','image-right','custom'],\n"
    '      "title": string,\n'
    '      "bullets": [string, ...] (optional),\n'
    '      "notes": string (optional),\n'
    '      "reuse_image_tag": string (optional)\n'
    "    }, ...\n"
    "  ]\n"
    "}. Keep bullets concise; avoid markdown formatting."
)

USER_GUIDANCE_SUFFIX = (
    "If the user guidance suggests a deck type (e.g., investor pitch, academic talk), structure appropriately: "
    "start with a title/summary, add sections, keep each slide focused. Limit to ~12–18 slides by default."
)

def call_openai(api_key: str, content: str, guidance: str) -> Dict[str, Any]:
    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": f"GUIDANCE: {guidance}\n\nTEXT:\n{content}\n\n{USER_GUIDANCE_SUFFIX}"},
    ]
    payload = {
        "model": LLM_MODELS["OpenAI"],
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
        "messages": messages,
    }
    r = session.post(url, headers=headers, json=payload, timeout=90)
    if r.status_code >= 400:
        raise RuntimeError(f"OpenAI error {r.status_code}: {r.text[:200]}")
    txt = r.json()["choices"][0]["message"]["content"]
    return extract_json(txt)

def call_openai_compatible(base_url: str, model: str, api_key: str, content: str, guidance: str) -> Dict[str, Any]:
    url = base_url.rstrip("/") + "/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": f"GUIDANCE: {guidance}\n\nTEXT:\n{content}\n\n{USER_GUIDANCE_SUFFIX}"},
    ]
    payload = {"model": model, "temperature": 0.2, "response_format": {"type": "json_object"}, "messages": messages}
    r = session.post(url, headers=headers, json=payload, timeout=90)
    if r.status_code >= 400:
        raise RuntimeError(f"OpenAI-compatible error {r.status_code}: {r.text[:200]}")
    data = r.json()
    # Try standard shape
    txt = data.get("choices", [{}])[0].get("message", {}).get("content") or ""
    if not txt:
        # Some providers return 'content' as list/array or different keys
        txt = json.dumps(data)
    return extract_json(txt)

def call_anthropic(api_key: str, content: str, guidance: str) -> Dict[str, Any]:
    url = "https://api.anthropic.com/v1/messages"
    headers = {"x-api-key": api_key, "anthropic-version": "2023-06-01", "content-type": "application/json"}
    prompt = SYSTEM_PROMPT + "\n\n" + f"GUIDANCE: {guidance}\n\nTEXT:\n{content}\n\n{USER_GUIDANCE_SUFFIX}"
    payload = {"model": LLM_MODELS["Anthropic"], "max_tokens": 4000, "messages": [{"role": "user", "content": prompt}], "temperature": 0.2}
    r = session.post(url, headers=headers, json=payload, timeout=90)
    if r.status_code >= 400:
        raise RuntimeError(f"Anthropic error {r.status_code}: {r.text[:200]}")
    parts = r.json().get("content", [])
    txt = "".join([p.get("text", "") for p in parts if p.get("type") == "text"])
    return extract_json(txt)

def call_gemini(api_key: str, content: str, guidance: str) -> Dict[str, Any]:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{LLM_MODELS['Gemini']}:generateContent?key={api_key}"
    prompt = SYSTEM_PROMPT + "\n\n" + f"GUIDANCE: {guidance}\n\nTEXT:\n{content}\n\n{USER_GUIDANCE_SUFFIX}"
    payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}], "generationConfig": {"temperature": 0.2, "responseMimeType": "application/json"}}
    r = session.post(url, json=payload, timeout=90)
    if r.status_code >= 400:
        raise RuntimeError(f"Gemini error {r.status_code}: {r.text[:200]}")
    candidates = r.json().get("candidates", [])
    txt = ""
    if candidates:
        parts = candidates[0].get("content", {}).get("parts", [])
        txt = "".join(p.get("text", "") for p in parts)
    return extract_json(txt or json.dumps({"slides": []}))

def call_openrouter(api_key: str, content: str, guidance: str) -> Dict[str, Any]:
    # OpenRouter generally supports OpenAI-compatible /chat/completions
    base = "https://openrouter.ai/api"
    return call_openai_compatible(base, LLM_MODELS["OpenRouter"], api_key, content, guidance)

PROVIDERS = {
    "OpenAI": call_openai,
    "Anthropic": call_anthropic,
    "Gemini": call_gemini,
    "OpenRouter": call_openrouter,
    # Custom handled separately in UI
}

# --------------------- PPTX Helpers ---------------------

def list_template_images(pptx_path: Path) -> List[Tuple[str, bytes]]:
    """Return [(filename, data), ...] from ppt/media of the uploaded template."""
    images: List[Tuple[str, bytes]] = []
    with zipfile.ZipFile(pptx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith('ppt/media/') and name.lower().split('.')[-1] in {"png", "jpg", "jpeg"}:
                images.append((Path(name).name, z.read(name)))
    return images

def find_layout(pres: Presentation, hint: str):
    """Best-effort mapping from layout_hint → a real layout in the template."""
    hint = (hint or "").lower()
    preferred_order = []
    if hint in {"title", "title slide"}:
        preferred_order = ["title slide", "title"]
    elif hint in {"section", "section header"}:
        preferred_order = ["section header", "title slide", "title and content"]
    elif hint in {"two-content", "two content", "comparison"}:
        preferred_order = ["two content", "comparison", "title and content"]
    elif hint in {"quote"}:
        preferred_order = ["quote", "title only", "blank", "title and content"]
    elif hint in {"image-left", "image-right"}:
        preferred_order = ["picture with caption", "title and content", "blank"]
    elif hint in {"bullets-only", "title+content", "title and content", "custom"}:
        preferred_order = ["title and content", "title only", "blank"]
    else:
        preferred_order = ["title and content", "title slide", "section header", "two content", "comparison", "title only", "blank"]

    def norm(n: str) -> str:
        return (n or "").strip().lower()

    for target in preferred_order:
        for layout in pres.slide_layouts:
            if target in norm(layout.name):
                return layout
    return pres.slide_layouts[0]

def add_text_to_placeholder(shape, text: str, is_title=False):
    if not shape:
        return
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (text or "").strip()
    if is_title:
        p.level = 0
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

def add_bullets_to_placeholder(shape, bullets: List[str]):
    if not shape or not bullets:
        return
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = (b or "").strip()
        p.level = 0

def place_image_on_slide(slide, img_bytes: bytes, placeholder_like=None):
    """Place image using dimensions from first content placeholder if available; else center it."""
    if placeholder_like is not None:
        left = placeholder_like.left
        top = placeholder_like.top
        width = placeholder_like.width
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)
            return
        except Exception:
            pass
    # Fallback: centered medium
    slide_width = slide.part.presentation.slide_width
    slide_height = slide.part.presentation.slide_height
    width = int(slide_width * 0.5)
    left = int((slide_width - width) / 2)
    top = int(slide_height * 0.25)
    slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)

def first_title_placeholder(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
            return shape
    return None

def first_content_placeholder(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and getattr(shape, "has_text_frame", False) and shape.placeholder_format.type in {2, 7}:  # BODY/CONTENT
            return shape
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            return shape
    return None

# --------------------- App UI ---------------------

st.set_page_config(page_title=APP_TITLE, layout="wide")
st.title(APP_TITLE)
st.caption("Paste text → choose provider → upload PowerPoint template → get a deck that matches the template style. No AI images.")

with st.sidebar:
    st.header("Settings")
    provider = st.selectbox("LLM Provider", ["OpenAI", "Anthropic", "Gemini", "OpenRouter", "Custom (OpenAI-compatible)"], index=0)
    api_key = st.text_input(f"{provider} API Key", type="password", help="Keys are used only for this session and never stored.")
    custom_base = ""
    custom_model = ""
    if provider == "Custom (OpenAI-compatible)":
        custom_base = st.text_input("Base URL (OpenAI-compatible)", value="https://api.openai.com")
        custom_model = st.text_input("Model name", value=LLM_MODELS[provider])
    st.markdown(":lock: The app never logs or stores API keys.")

col1, col2 = st.columns([2, 1])
with col1:
    text = st.text_area("Paste text / markdown / long-form prose", height=300, placeholder="Paste your content here…")
    presets = {
        "None": "",
        "Investor Pitch": "Turn into an investor pitch deck with problem-solution, market, traction, business model, team, ask.",
        "Sales Deck": "Create a persuasive sales deck with problem, product, benefits, proof, pricing, next steps.",
        "Research Summary": "Summarize as an academic research talk with abstract, background, method, results, discussion, conclusion.",
        "Technical Design Review": "Convert into a technical design review with context, requirements, architecture, trade-offs, risks, timeline.",
        "Visual-Heavy Teaser": "Make a minimal-text, visual-heavy teaser with big takeaways and short bullets.",
    }
    preset = st.selectbox("Use-case guidance template (optional)", list(presets.keys()), index=0)
    guidance = st.text_input("Optional one-line guidance (tone/use case)", value=presets.get(preset, ""))

with col2:
    template_file = st.file_uploader("Upload a PowerPoint template/presentation (.pptx or .potx)", type=["pptx", "potx"], accept_multiple_files=False)
    max_slides = st.slider("Max slides", 5, 40, 16)
    include_notes = st.checkbox("Auto-generate speaker notes (if provided by LLM)", value=True)
    st.info("Tip: Keep uploads ≤ %d MB. Images are reused from the template; no AI images." % MAX_UPLOAD_MB)

def plan_slides() -> Dict[str, Any]:
    if not api_key:
        st.error("Please enter your API key in the sidebar.")
        return {}
    if not text or len(text.strip()) < 5:
        st.error("Please paste sufficient input text.")
        return {}
    if provider == "Custom (OpenAI-compatible)" and not custom_base:
        st.error("Please provide the base URL for your custom provider.")
        return {}

    try:
        with st.spinner("Analyzing text and planning slides…"):
            if provider == "Custom (OpenAI-compatible)":
                plan = call_openai_compatible(custom_base, custom_model or LLM_MODELS[provider], api_key, text, guidance or "")
            else:
                call = PROVIDERS.get(provider)
                plan = call(api_key, text, guidance or "")
    except Exception as e:
        st.error(f"LLM call failed: {e}")
        return {}

    slides = plan.get("slides", [])
    if not slides:
        slides = [{"layout_hint": "title", "title": "Overview", "bullets": [text[:1000]]}]
    if len(slides) > max_slides:
        slides = slides[:max_slides]
    plan["slides"] = slides
    return plan

def build_pptx(plan: Dict[str, Any], template_bytes: bytes) -> bytes:
    pres = Presentation(io.BytesIO(template_bytes))

    # Save template temporarily to extract images
    tmp_path = Path(".") / "_user_template_tmp.pptx"
    tmp_path.write_bytes(template_bytes)
    images = list_template_images(tmp_path)

    deck_title = plan.get("title")
    if deck_title:
        layout = find_layout(pres, "title")
        slide = pres.slides.add_slide(layout)
        tph = first_title_placeholder(slide)
        if tph:
            add_text_to_placeholder(tph, deck_title, is_title=True)

    img_cycle = images.copy()

    for s in plan.get("slides", []):
        layout = find_layout(pres, s.get("layout_hint") or "title+content")
        slide = pres.slides.add_slide(layout)

        title_shape = first_title_placeholder(slide)
        body_shape = first_content_placeholder(slide)

        add_text_to_placeholder(title_shape, s.get("title", ""), is_title=True)
        bullets = s.get("bullets") or []
        add_bullets_to_placeholder(body_shape, bullets)

        if s.get("notes") and include_notes:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = s.get("notes", "")

        wants_image = (s.get("layout_hint") in {"image-left", "image-right", "quote", "title+content"}) or bool(s.get("reuse_image_tag"))
        if wants_image and img_cycle:
            fname, data = img_cycle[0]
            img_cycle = img_cycle[1:] + [(fname, data)]
            place_image_on_slide(slide, data, placeholder_like=body_shape)

    out = io.BytesIO()
    pres.save(out)
    return out.getvalue()

plan_holder = st.empty()
export_holder = st.empty()

if st.button("Generate deck", type="primary"):
    if not template_file:
        st.error("Please upload a .pptx/.potx template to carry over style and images.")
    else:
        if template_file.size > MAX_UPLOAD_MB * 1024 * 1024:
            st.error(f"Template file too large (>{MAX_UPLOAD_MB} MB). Please upload a smaller file.")
        else:
            plan = plan_slides()
            if plan:
                # Preview table
                with st.expander("Slide plan preview"):
                    rows = []
                    for i, s in enumerate(plan.get("slides", []), start=1):
                        rows.append({
                            "No.": i,
                            "Title": s.get("title", "")[:80],
                            "Bullets": len(s.get("bullets") or []),
                            "Layout hint": s.get("layout_hint", ""),
                        })
                    try:
                        import pandas as pd
                        st.dataframe(pd.DataFrame(rows))
                    except Exception:
                        st.write(rows)
                    st.code(json.dumps(plan, indent=2)[:100000])

                with st.spinner("Building PowerPoint…"):
                    pptx_bytes = build_pptx(plan, template_file.getvalue())
                st.success("Presentation ready!")
                st.download_button(
                    label="Download .pptx",
                    data=pptx_bytes,
                    file_name="generated_deck.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

st.markdown("---")
with st.expander("How it works & guarantees"):
    st.markdown(
        """
        - **Template-aware**: New slides are added to your uploaded template, inheriting theme, colors, fonts, and layouts.
        - **Image reuse only**: The app never generates images. If the template includes images (`ppt/media/`), they'll be reused.
        - **Your API key**: You provide an OpenAI / Anthropic / Gemini / OpenRouter key, or a custom OpenAI-compatible base URL. Keys are **not stored**.
        - **Structured output**: The LLM returns a JSON slide plan. If parsing fails, a safe fallback is used.
        - **Notes**: If present and enabled, speaker notes are included.
        - **Limits**: Reasonable file size limits and API retries are enabled.
        """
    )
